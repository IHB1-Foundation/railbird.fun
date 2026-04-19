#!/usr/bin/env python3
"""Generate Railbird hackathon pitch deck (PPTX).

Brand identity sourced from apps/web/src/app/globals.css.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Railbird Brand Palette (from globals.css) ────────────────
BG_DEEP     = RGBColor(0x06, 0x07, 0x0B)   # --background
BG_SOFT     = RGBColor(0x0D, 0x10, 0x20)   # --background-soft
FG          = RGBColor(0xF7, 0xF8, 0xFF)   # --foreground
ACCENT      = RGBColor(0x81, 0x6C, 0xF9)   # --accent (primary purple)
ACCENT_SOFT = RGBColor(0xC8, 0xA9, 0xFF)   # --accent-soft
NEON_VIOLET = RGBColor(0xA7, 0x8B, 0xFA)   # --neon-violet
NEON_CYAN   = RGBColor(0x22, 0xD3, 0xEE)   # --neon-cyan
NEON_GOLD   = RGBColor(0xFA, 0xCC, 0x15)   # --neon-gold
NEON_LIME   = RGBColor(0xA3, 0xE6, 0x35)   # --success
DANGER      = RGBColor(0xEF, 0x44, 0x44)   # --danger
MUTED       = RGBColor(0x7C, 0x84, 0xA2)   # --text-muted
MUTED_SOFT  = RGBColor(0x9B, 0xA3, 0xC1)   # --text-muted-soft
CARD_BG     = RGBColor(0x10, 0x12, 0x21)   # approx --card-bg
CARD_BORDER = RGBColor(0x23, 0x27, 0x3A)   # approx --card-border
GOLD_TEXT   = RGBColor(0xFF, 0xBF, 0x4D)   # --text-gold

# Fonts matching product
FONT_DISPLAY = "Calibri"    # closest to Space Grotesk in system fonts
FONT_BODY    = "Calibri"    # closest to Roboto
FONT_MONO    = "Consolas"   # closest to JetBrains Mono

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# ── Layout margins ────────────────────────────────────────────
L_MARGIN = Inches(1.2)
R_MARGIN = Inches(1.2)
CONTENT_W = SLIDE_W - L_MARGIN - R_MARGIN


# ── Helpers ───────────────────────────────────────────────────

def set_slide_bg(slide, color=BG_DEEP):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _configure_run(run, font_size, color, bold, font_name):
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = font_name


def make_textbox(slide, left, top, width, height, *,
                 anchor=MSO_ANCHOR.TOP, word_wrap=True,
                 margin_left=Inches(0), margin_right=Inches(0),
                 margin_top=Inches(0), margin_bottom=Inches(0)):
    """Create a text box and return the text frame, configured for no-overlap."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = word_wrap
    tf.auto_size = None  # fixed-size frame — text won't spill
    tf.margin_left = margin_left
    tf.margin_right = margin_right
    tf.margin_top = margin_top
    tf.margin_bottom = margin_bottom
    return tf


def write_para(tf, text, size=18, color=FG, bold=False, font=FONT_BODY,
               align=PP_ALIGN.LEFT, space_before=Pt(0), space_after=Pt(0),
               first=False):
    """Write a paragraph to a text frame. If first=True, reuse the default paragraph."""
    if first:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = text
    p.alignment = align
    p.space_before = space_before
    p.space_after = space_after
    _configure_run(p.runs[0] if p.runs else p, size, color, bold, font)
    # Ensure the run exists
    if not p.runs:
        run = p.add_run()
        run.text = text
        p.text = ""
        # re-add via run
        _configure_run(run, size, color, bold, font)
    return p


def add_rect(slide, left, top, w, h, fill_color, border_color=None,
             border_width=Pt(0), radius=True):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(shape_type, left, top, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill_color
    if border_color:
        s.line.color.rgb = border_color
        s.line.width = border_width
    else:
        s.line.fill.background()
    return s


def add_thin_line(slide, left, top, width, color=ACCENT):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(3))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()


def slide_header(slide, tag, title, tag_color=ACCENT):
    """Add a consistent slide header: small tag + thin line + big title."""
    # Tag label
    tf = make_textbox(slide, L_MARGIN, Inches(0.6), CONTENT_W, Inches(0.35))
    write_para(tf, tag.upper(), size=12, color=tag_color, bold=True,
               font=FONT_MONO, first=True)
    # Accent line
    add_thin_line(slide, L_MARGIN, Inches(0.95), Inches(1.2), tag_color)
    # Title
    tf = make_textbox(slide, L_MARGIN, Inches(1.15), CONTENT_W, Inches(0.7))
    write_para(tf, title, size=30, color=FG, bold=True, font=FONT_DISPLAY,
               first=True)


def card_with_text(slide, left, top, w, h, title, bullets, *,
                   title_color=ACCENT, border_color=CARD_BORDER,
                   title_size=18, bullet_size=13):
    """A card shape with internal text frame — no separate text boxes."""
    shape = add_rect(slide, left, top, w, h, CARD_BG,
                     border_color=border_color, border_width=Pt(1))

    # Use shape's own text frame for the title + bullets
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.margin_left = Inches(0.25)
    tf.margin_right = Inches(0.25)
    tf.margin_top = Inches(0.2)
    tf.margin_bottom = Inches(0.15)

    # Title paragraph
    p = tf.paragraphs[0]
    p.text = title
    p.space_after = Pt(8)
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    _configure_run(run, title_size, title_color, True, FONT_DISPLAY)

    # Bullet paragraphs
    for line in bullets:
        p = tf.add_paragraph()
        p.text = line
        p.space_before = Pt(3)
        p.space_after = Pt(1)
        p.alignment = PP_ALIGN.LEFT
        run = p.runs[0]
        _configure_run(run, bullet_size, MUTED_SOFT, False, FONT_BODY)

    return shape


# ═══════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s)

# Decorative left bar
add_rect(s, Inches(0), Inches(0), Inches(0.08), SLIDE_H, ACCENT,
         radius=False)

# Vertical accent line
add_thin_line(s, Inches(1.6), Inches(2.1), Inches(2.5), ACCENT)

# Title
tf = make_textbox(s, Inches(1.6), Inches(2.3), Inches(9), Inches(1))
write_para(tf, "RAILBIRD", size=60, color=FG, bold=True,
           font=FONT_DISPLAY, first=True)

# Subtitle
tf = make_textbox(s, Inches(1.6), Inches(3.4), Inches(9), Inches(0.7))
write_para(tf, "The World's First Trustless AI Poker Protocol",
           size=24, color=ACCENT_SOFT, bold=False, font=FONT_DISPLAY,
           first=True)

# Context line
tf = make_textbox(s, Inches(1.6), Inches(4.5), Inches(9), Inches(0.5))
write_para(tf, "On-Chain Horizon Hackathon  \u2502  AI Track  \u2502  HashKey Chain",
           size=14, color=MUTED, font=FONT_BODY, first=True)

# URL
tf = make_textbox(s, Inches(1.6), Inches(5.3), Inches(4), Inches(0.4))
write_para(tf, "railbird.fun", size=18, color=NEON_CYAN, bold=True,
           font=FONT_MONO, first=True)


# ═══════════════════════════════════════════════════════════════
# SLIDE 2 — Problem
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s)
slide_header(s, "Problem", "On-chain gaming has two trust failures", tag_color=DANGER)

# Left column — problem 1
card_with_text(s, L_MARGIN, Inches(2.2), Inches(5.2), Inches(3),
               "Centralized Dealers",
               ["\u2022  The house sees every card before you do",
                "\u2022  Outcomes can be silently manipulated",
                "\u2022  Zero on-chain proof of fairness",
                "\u2022  Players must blindly trust the operator"],
               title_color=DANGER, border_color=RGBColor(0x3D, 0x1A, 0x1A))

# Right column — problem 2
card_with_text(s, Inches(7), Inches(2.2), Inches(5.2), Inches(3),
               "Opaque AI Agents",
               ["\u2022  AI makes DeFi decisions with zero transparency",
                "\u2022  No way to verify what the AI actually decided",
                "\u2022  Black-box strategies erode user trust",
                "\u2022  No auditability after the fact"],
               title_color=DANGER, border_color=RGBColor(0x3D, 0x1A, 0x1A))

# Bottom question
tf = make_textbox(s, L_MARGIN, Inches(5.7), CONTENT_W, Inches(0.7))
write_para(tf, "What if every shuffle was provably fair\nand every AI decision was recorded on-chain?",
           size=17, color=ACCENT_SOFT, bold=True, font=FONT_DISPLAY,
           align=PP_ALIGN.LEFT, first=True)


# ═══════════════════════════════════════════════════════════════
# SLIDE 3 — Solution (3-column, varied accent)
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s)
slide_header(s, "Solution", "Fully on-chain poker with verifiable AI")

COL_W = Inches(3.5)
COL_GAP = Inches(0.3)
col_start = L_MARGIN

pillars = [
    ("Trustless Dealer", NEON_LIME,
     ["\u2022  VRF randomness + dealer pre-commit seed",
      "\u2022  Deterministic Fisher-Yates shuffle",
      "\u2022  Verified on-chain at every showdown",
      "\u2022  No entity can see or manipulate the deck"]),
    ("ECIES Encrypted Cards", NEON_CYAN,
     ["\u2022  Wallet public key encrypts each player's cards",
      "\u2022  Only the seat owner can decrypt",
      "\u2022  keccak256 commit / reveal integrity",
      "\u2022  Not even our servers can read your hand"]),
    ("Verifiable AI", NEON_VIOLET,
     ["\u2022  Gemini 2.0 Flash reasoning engine",
      "\u2022  Decision hash committed on-chain",
      "\u2022  Anyone can audit any action, any time",
      "\u2022  Full transparency without trust"]),
]

for i, (title, color, bullets) in enumerate(pillars):
    left = col_start + i * (COL_W + COL_GAP)
    # Color bar on top of card
    add_rect(s, left, Inches(2.2), COL_W, Pt(4), color, radius=False)
    card_with_text(s, left, Inches(2.3), COL_W, Inches(3.2),
                   title, bullets,
                   title_color=color, border_color=CARD_BORDER)


# ═══════════════════════════════════════════════════════════════
# SLIDE 4 — Architecture (layered, not 3 identical cards)
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s)
slide_header(s, "Architecture", "Three layers on HashKey Chain")

# Full-width on-chain bar
add_rect(s, L_MARGIN, Inches(2.2), CONTENT_W, Inches(0.45), ACCENT,
         radius=False)
tf = make_textbox(s, L_MARGIN, Inches(2.2), CONTENT_W, Inches(0.45),
                  margin_left=Inches(0.2))
write_para(tf, "ON-CHAIN  \u2014  HashKey Chain Testnet (ID: 133)", size=13,
           color=BG_DEEP, bold=True, font=FONT_MONO, first=True)

# On-chain contracts row
contracts = ["PokerTable", "SideBetPool", "VRFAdapter", "ChipToken",
             "PlayerRegistry", "PlayerVault"]
cw = Inches(1.7)
for i, name in enumerate(contracts):
    left = L_MARGIN + Inches(0.05) + i * (cw + Inches(0.12))
    shape = add_rect(s, left, Inches(2.8), cw, Inches(0.55), CARD_BG,
                     border_color=ACCENT, border_width=Pt(1))
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.margin_left = Inches(0.08)
    tf.margin_top = Inches(0.05)
    p = tf.paragraphs[0]
    p.text = name
    p.alignment = PP_ALIGN.CENTER
    _configure_run(p.runs[0], 11, ACCENT_SOFT, True, FONT_MONO)

# Off-chain services band
add_rect(s, L_MARGIN, Inches(3.7), CONTENT_W, Inches(0.4), NEON_CYAN,
         radius=False)
tf = make_textbox(s, L_MARGIN, Inches(3.7), CONTENT_W, Inches(0.4),
                  margin_left=Inches(0.2))
write_para(tf, "OFF-CHAIN SERVICES", size=13, color=BG_DEEP, bold=True,
           font=FONT_MONO, first=True)

# Off-chain detail cards (2 wide)
card_with_text(s, L_MARGIN, Inches(4.2), Inches(5.2), Inches(2),
               "Indexer + WebSocket API",
               ["Contract events \u2192 Postgres \u2192 REST + WS",
                "Real-time table state streaming",
                "Leaderboard aggregation"],
               title_color=NEON_CYAN)

card_with_text(s, Inches(7), Inches(4.2), Inches(5.2), Inches(2),
               "OwnerView + Dealer Service",
               ["Wallet-signature authentication",
                "ECIES hole card delivery (owner-only ACL)",
                "KeeperBot for liveness (forceTimeout)"],
               title_color=NEON_CYAN)

# AI band
add_rect(s, L_MARGIN, Inches(6.5), CONTENT_W, Inches(0.4), NEON_VIOLET,
         radius=False)
tf = make_textbox(s, L_MARGIN, Inches(6.5), CONTENT_W, Inches(0.4),
                  margin_left=Inches(0.2))
write_para(tf, "AI LAYER  \u2014  4 Gemini 2.0 Flash agents  \u2502  hand strength + pot odds + opponent modeling  \u2502  on-chain decision hash",
           size=12, color=BG_DEEP, bold=True, font=FONT_MONO, first=True)


# ═══════════════════════════════════════════════════════════════
# SLIDE 5 — AI Agents (horizontal layout with colored borders)
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s)
slide_header(s, "AI Agents", "Four Gemini-powered personalities")

agents = [
    ("Aegis",    "Tight",     "0.2",
     "Patient, disciplined.\nWaits for premium hands.\nRarely bluffs.", NEON_CYAN),
    ("Maverick", "Balanced",  "0.4",
     "Reads opponents, adapts.\nMixes value bets and bluffs.\nSolid fundamentals.", ACCENT),
    ("Nova",     "Loose",     "0.6",
     "Creative lines.\nPlays many hands.\nFinds unconventional spots.", NEON_GOLD),
    ("Rex",      "Maniac",    "0.8",
     "Maximum pressure.\nRelentless aggression.\nForces hard decisions.", DANGER),
]

card_w = Inches(2.6)
card_h = Inches(4)
gap = Inches(0.3)
start_x = L_MARGIN + Inches(0.1)

for i, (name, style, aggr, desc, color) in enumerate(agents):
    left = start_x + i * (card_w + gap)
    top = Inches(2.2)

    # Card background
    shape = add_rect(s, left, top, card_w, card_h, CARD_BG,
                     border_color=color, border_width=Pt(2))
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.35)
    tf.margin_bottom = Inches(0.15)

    # Name
    p = tf.paragraphs[0]
    p.text = name
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(4)
    _configure_run(p.runs[0], 26, color, True, FONT_DISPLAY)

    # Style badge
    p = tf.add_paragraph()
    p.text = style
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(12)
    _configure_run(p.runs[0], 14, FG, True, FONT_BODY)

    # Aggression line
    p = tf.add_paragraph()
    p.text = f"Aggression  {aggr}"
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(14)
    _configure_run(p.runs[0], 11, MUTED, False, FONT_MONO)

    # Description lines
    for line in desc.split("\n"):
        p = tf.add_paragraph()
        p.text = line
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(2)
        _configure_run(p.runs[0], 12, MUTED_SOFT, False, FONT_BODY)

    # Colored top bar
    add_rect(s, left, top, card_w, Pt(4), color, radius=False)


# ═══════════════════════════════════════════════════════════════
# SLIDE 6 — HashKey Chain (2x2 grid instead of 4 narrow cards)
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s)
slide_header(s, "HashKey Chain", "Why we chose this chain", tag_color=NEON_CYAN)

features = [
    ("Wallet-Based Identity", NEON_VIOLET,
     ["All auth via wallet signatures",
      "No email/password accounts needed",
      "On-chain ownership = authorization"]),
    ("OP Stack EVM Equivalence", NEON_LIME,
     ["Standard Solidity + Foundry tooling",
      "Zero contract modifications needed",
      "Low gas costs, familiar dev experience"]),
    ("VRF On-Chain Randomness", NEON_CYAN,
     ["Backbone of the trustless dealer",
      "Provably fair community card shuffles",
      "Verifiable on-chain at showdown"]),
    ("Blockscout Explorer", NEON_GOLD,
     ["All 6 contracts source-verified",
      "Public code inspection by anyone",
      "Full on-chain state transparency"]),
]

half_w = Inches(5.2)
half_h = Inches(2)
gap_x = Inches(0.5)
gap_y = Inches(0.35)

for i, (title, color, bullets) in enumerate(features):
    col = i % 2
    row = i // 2
    left = L_MARGIN + col * (half_w + gap_x)
    top = Inches(2.2) + row * (half_h + gap_y)

    # Colored left bar
    add_rect(s, left, top, Pt(4), half_h, color, radius=False)
    card_with_text(s, left, top, half_w, half_h,
                   title, [f"\u2022  {b}" for b in bullets],
                   title_color=color, title_size=17, bullet_size=13)

# Footer
tf = make_textbox(s, L_MARGIN, Inches(6.6), CONTENT_W, Inches(0.4))
write_para(tf, "6 contracts deployed & source-verified on HashKey Chain Testnet (Chain ID: 133)",
           size=13, color=MUTED, font=FONT_MONO, align=PP_ALIGN.CENTER,
           first=True)


# ═══════════════════════════════════════════════════════════════
# SLIDE 7 — Live Demo (full-width stacked sections)
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s)
slide_header(s, "Live Demo", "railbird.fun")

# Large URL
tf = make_textbox(s, L_MARGIN, Inches(2.0), CONTENT_W, Inches(0.5))
write_para(tf, "Everything on-chain. Everything verifiable.",
           size=18, color=ACCENT_SOFT, font=FONT_DISPLAY, first=True)

demo_items = [
    ("\u2660  Live Table Viewer",
     "Real-time community cards, pot, chip stacks  \u2502  Action log with block numbers  \u2502  VRF status widget"),
    ("\U0001F916  AI Decision Transparency",
     "Hand strength + pot odds breakdown per action  \u2502  On-chain reasoning hash verification  \u2502  Confidence gauge"),
    ("\u2728  Showdown & Settlement",
     "Card flip animation  \u2502  Winner highlight + pot distribution  \u2502  Commit/reveal verified on-chain"),
    ("\U0001F3C6  Leaderboard & Agent Pages",
     "ROI / PnL / Win Rate / MDD rankings  \u2502  Vault metrics (A/B/N/P)  \u2502  Token trading widget"),
]

row_h = Inches(0.9)
for i, (title, desc) in enumerate(demo_items):
    top = Inches(2.8) + i * (row_h + Inches(0.15))
    shape = add_rect(s, L_MARGIN, top, CONTENT_W, row_h, CARD_BG,
                     border_color=CARD_BORDER, border_width=Pt(1))
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.margin_left = Inches(0.3)
    tf.margin_right = Inches(0.3)
    tf.margin_top = Inches(0.12)

    p = tf.paragraphs[0]
    p.text = title
    p.space_after = Pt(4)
    _configure_run(p.runs[0], 16, ACCENT, True, FONT_DISPLAY)

    p = tf.add_paragraph()
    p.text = desc
    _configure_run(p.runs[0], 12, MUTED_SOFT, False, FONT_BODY)


# ═══════════════════════════════════════════════════════════════
# SLIDE 8 — Sidebet Market (the economic engine)
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s)
slide_header(s, "AI Prediction Market", "Spectating becomes evaluating", tag_color=NEON_GOLD)

# ── Left half: How it works ──
tf = make_textbox(s, L_MARGIN, Inches(2.1), Inches(5.5), Inches(0.4))
write_para(tf, "HOW IT WORKS", size=12, color=NEON_GOLD, bold=True,
           font=FONT_MONO, first=True)

flow_steps = [
    ("1.  Watch", "Spectators watch AI agents play poker in real-time."),
    ("2.  Bet", "Pick which agent wins the hand. Place RCHIP on that seat."),
    ("3.  Settle", "Hand settles on-chain. SideBetPool reads the winner from PokerTable."),
    ("4.  Claim", "Winners claim proportional payout: (your bet \u00d7 total pool) / winner-seat total."),
]

for i, (step, desc) in enumerate(flow_steps):
    top = Inches(2.7) + i * Inches(1.0)
    # Step number
    tf = make_textbox(s, L_MARGIN, top, Inches(1.2), Inches(0.4))
    write_para(tf, step, size=15, color=NEON_GOLD, bold=True,
               font=FONT_DISPLAY, first=True)
    # Description
    tf = make_textbox(s, L_MARGIN + Inches(1.3), top, Inches(4.2), Inches(0.8))
    write_para(tf, desc, size=13, color=MUTED_SOFT, font=FONT_BODY,
               first=True)

# ── Right half: Why it matters (DeFi angle) ──
right_x = Inches(7)
tf = make_textbox(s, right_x, Inches(2.1), Inches(5.5), Inches(0.4))
write_para(tf, "WHY IT MATTERS FOR AI", size=12, color=ACCENT, bold=True,
           font=FONT_MONO, first=True)

defi_points = [
    ("AI Evaluation Layer", ACCENT,
     "Spectators actively assess which AI strategy performs best. Crowd-sourced AI evaluation with real stakes."),
    ("Transparent & On-Chain", NEON_CYAN,
     "Pari-mutuel, fully on-chain. No house edge. Dynamic odds from agent skill profiles \u2014 purely proportional."),
    ("Open Platform", NEON_LIME,
     "SideBetPool.sol is permissionless. Anyone can build prediction interfaces, analysis bots, or trackers."),
    ("Auto Settlement", NEON_VIOLET,
     "KeeperBot auto-settles after showdown. Winner recorded on-chain. Claims are permissionless."),
]

for i, (title, color, desc) in enumerate(defi_points):
    top = Inches(2.7) + i * Inches(1.0)
    add_rect(s, right_x, top + Inches(0.05), Inches(0.1), Inches(0.1),
             color, radius=False)
    tf = make_textbox(s, right_x + Inches(0.2), top, Inches(5.2), Inches(0.85))
    write_para(tf, title, size=14, color=color, bold=True,
               font=FONT_DISPLAY, first=True)
    write_para(tf, desc, size=11, color=MUTED_SOFT, font=FONT_BODY,
               space_before=Pt(3))


# ═══════════════════════════════════════════════════════════════
# SLIDE 9 — Security
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s)
slide_header(s, "Security", "Trust model & on-chain guarantees")

sec_items = [
    ("Commit / Reveal", NEON_LIME,
     "keccak256 commitments stored on-chain. At showdown, cards + salt are revealed and the contract verifies the match. Post-hoc integrity is guaranteed."),
    ("Anti-MEV", NEON_CYAN,
     "One action per block per table. Prevents front-running and sandwich attacks. Deterministic ordering is enforced at the contract level."),
    ("Liveness Guarantees", NEON_VIOLET,
     "30-minute turn timeouts. Any address can call forceTimeout() with keeper incentives. The game always moves forward."),
    ("Non-Dilutive Treasury", NEON_GOLD,
     "Treasury buys tokens only below NAV/share, sells only above. Smart contract reverts if the constraint is violated. Existing holders are never diluted."),
]

for i, (title, color, desc) in enumerate(sec_items):
    top = Inches(2.2) + i * Inches(1.2)
    # Color dot
    add_rect(s, L_MARGIN, top + Inches(0.08), Inches(0.12), Inches(0.12),
             color, radius=False)
    # Title + description in one text frame
    tf = make_textbox(s, L_MARGIN + Inches(0.3), top, CONTENT_W - Inches(0.3),
                      Inches(1.05))
    write_para(tf, title, size=16, color=color, bold=True,
               font=FONT_DISPLAY, first=True)
    write_para(tf, desc, size=12, color=MUTED_SOFT, font=FONT_BODY,
               space_before=Pt(4))


# ═══════════════════════════════════════════════════════════════
# SLIDE 10 — Ecosystem Impact (why this matters for HashKey Chain)
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s)
slide_header(s, "Ecosystem Impact", "An open AI arena on HashKey Chain", tag_color=NEON_CYAN)

# ── Left: flywheel narrative ──
tf = make_textbox(s, L_MARGIN, Inches(2.1), Inches(5.5), Inches(0.4))
write_para(tf, "ACTIVITY FLYWHEEL", size=12, color=NEON_CYAN, bold=True,
           font=FONT_MONO, first=True)

flywheel = [
    ("Open AI Arena",
     "Anyone can deploy custom AI agents via web wizard. "
     "Different strategies compete, adapt, and evolve \u2014 all on-chain, all verifiable."),
    ("Natural Onboarding",
     "Spectating is free, no wallet. Predicting outcomes or deploying agents "
     "converts passive viewers into active on-chain users."),
    ("Composable AI Infra",
     "Agent registry, prediction market, game protocol \u2014 all permissionless. "
     "Any team can build strategy analyzers, dashboards, or tournament platforms."),
    ("Continuous Activity",
     "AI agents play 24/7. Every action, prediction, and settlement "
     "is an on-chain transaction. Steady organic demand, not hype spikes."),
]

for i, (title, desc) in enumerate(flywheel):
    top = Inches(2.7) + i * Inches(1.1)
    tf = make_textbox(s, L_MARGIN, top, Inches(5.5), Inches(0.95))
    write_para(tf, title, size=15, color=NEON_CYAN, bold=True,
               font=FONT_DISPLAY, first=True)
    write_para(tf, desc, size=11, color=MUTED_SOFT, font=FONT_BODY,
               space_before=Pt(3))

# ── Right: concrete numbers / projections ──
right_x = Inches(7)
tf = make_textbox(s, right_x, Inches(2.1), Inches(5), Inches(0.4))
write_para(tf, "MAINNET PROJECTION", size=12, color=NEON_GOLD, bold=True,
           font=FONT_MONO, first=True)

# Metric cards
metrics = [
    ("20+",  "txs per hand",     "Poker actions + VRF + settlement + sidebet lifecycle", ACCENT),
    ("4\u00d7",   "multiplier",       "Each sidebet bettor adds bet + claim txs on top of base poker volume", NEON_GOLD),
    ("24/7", "autonomous",       "AI agents play continuously \u2014 no human scheduling, no downtime", NEON_LIME),
    ("0",    "cold start cost",  "Spectating is free, no wallet needed. Wallet connects only when betting", NEON_VIOLET),
]

for i, (num, label, desc, color) in enumerate(metrics):
    top = Inches(2.7) + i * Inches(1.1)
    # Big number
    tf = make_textbox(s, right_x, top, Inches(1.4), Inches(0.5))
    write_para(tf, num, size=28, color=color, bold=True,
               font=FONT_DISPLAY, align=PP_ALIGN.RIGHT, first=True)
    # Label + desc
    tf = make_textbox(s, right_x + Inches(1.6), top, Inches(3.8), Inches(0.9))
    write_para(tf, label, size=14, color=color, bold=True,
               font=FONT_BODY, first=True)
    write_para(tf, desc, size=11, color=MUTED, font=FONT_BODY,
               space_before=Pt(2))


# ═══════════════════════════════════════════════════════════════
# SLIDE 11 — Traction & Roadmap
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s)
slide_header(s, "Traction", "Built and deployed in 5 weeks")

# ── TODAY section (left half) ──
tf = make_textbox(s, L_MARGIN, Inches(2.2), Inches(5), Inches(0.35))
write_para(tf, "TODAY", size=12, color=ACCENT, bold=True,
           font=FONT_MONO, first=True)

stats = [
    ("6",    "Contracts Deployed", ACCENT),
    ("4",    "Autonomous AI Agents", NEON_VIOLET),
    ("100%", "Source Verified", NEON_LIME),
    ("Live", "railbird.fun", NEON_CYAN),
]

for i, (num, label, color) in enumerate(stats):
    top = Inches(2.7) + i * Inches(0.9)
    # Number
    tf = make_textbox(s, L_MARGIN, top, Inches(1.8), Inches(0.55))
    write_para(tf, num, size=32, color=color, bold=True,
               font=FONT_DISPLAY, align=PP_ALIGN.RIGHT, first=True)
    # Label
    tf = make_textbox(s, L_MARGIN + Inches(2.1), top + Inches(0.08),
                      Inches(3.5), Inches(0.45))
    write_para(tf, label, size=14, color=MUTED_SOFT, font=FONT_BODY,
               first=True)

# ── NEXT section (right half) ──
tf = make_textbox(s, Inches(7), Inches(2.2), Inches(5), Inches(0.35))
write_para(tf, "ROADMAP", size=12, color=NEON_GOLD, bold=True,
           font=FONT_MONO, first=True)

roadmap = [
    ("ZK Proofs",    "Fully trustless dealing \u2014 remove dealer trust entirely"),
    ("Tournaments",  "Multi-table brackets with progressive elimination"),
    ("Mobile",       "Optimized spectating UX with push notifications"),
    ("Mainnet",      "Real economic stakes on HashKey Chain mainnet"),
]

for i, (title, desc) in enumerate(roadmap):
    top = Inches(2.7) + i * Inches(0.9)
    tf = make_textbox(s, Inches(7), top, Inches(5.2), Inches(0.8))
    write_para(tf, title, size=16, color=NEON_GOLD, bold=True,
               font=FONT_DISPLAY, first=True)
    write_para(tf, desc, size=12, color=MUTED_SOFT, font=FONT_BODY,
               space_before=Pt(3))


# ═══════════════════════════════════════════════════════════════
# SLIDE 12 — Closing
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s)

# Left accent bar
add_rect(s, Inches(0), Inches(0), Inches(0.08), SLIDE_H, ACCENT,
         radius=False)

add_thin_line(s, Inches(1.6), Inches(2.3), Inches(2), ACCENT)

# Title
tf = make_textbox(s, Inches(1.6), Inches(2.5), Inches(9), Inches(0.9))
write_para(tf, "RAILBIRD", size=56, color=FG, bold=True,
           font=FONT_DISPLAY, first=True)

# Tagline
tf = make_textbox(s, Inches(1.6), Inches(3.4), Inches(9), Inches(0.6))
write_para(tf, "Trustless AI Poker. Fully On-Chain.",
           size=22, color=ACCENT_SOFT, font=FONT_DISPLAY, first=True)

# Three pillars
tf = make_textbox(s, Inches(1.6), Inches(4.4), Inches(9), Inches(2))
write_para(tf, "Every card provably fair.", size=16, color=NEON_LIME,
           bold=True, font=FONT_BODY, first=True, space_after=Pt(6))
write_para(tf, "Every AI decision verifiable.", size=16, color=NEON_CYAN,
           bold=True, font=FONT_BODY, space_after=Pt(6))
write_para(tf, "Every prediction settles on-chain.", size=16, color=NEON_VIOLET,
           bold=True, font=FONT_BODY, space_after=Pt(20))
write_para(tf, "railbird.fun", size=18, color=NEON_CYAN, bold=True,
           font=FONT_MONO, space_after=Pt(4))
write_para(tf, "Built on HashKey Chain", size=13, color=MUTED,
           font=FONT_BODY)


# ── Save ──────────────────────────────────────────────────────
out_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
out_path = os.path.join(out_dir, "Railbird_Pitch_Deck.pptx")
prs.save(out_path)
print(f"Deck saved: {out_path}")
print(f"Slides: {len(prs.slides)}")
